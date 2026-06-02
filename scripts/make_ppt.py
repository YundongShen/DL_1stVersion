"""
make_ppt.py — Generate EEL_presentation.pptx  (10-minute talk)
Run:  pip install python-pptx && python3 make_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette (light theme) ───────────────────────────────
BG          = RGBColor(0xFF, 0xFF, 0xFF)
DARK        = RGBColor(0x1A, 0x23, 0x7E)   # deep navy — headings
TEXT        = RGBColor(0x21, 0x21, 0x21)   # near black — body
ACCENT      = RGBColor(0x15, 0x65, 0xC0)   # deep blue
SUB         = RGBColor(0x45, 0x56, 0x70)   # steel gray — secondary
GREEN       = RGBColor(0x2E, 0x7D, 0x32)
AMBER       = RGBColor(0xE6, 0x51, 0x00)
RED         = RGBColor(0xC6, 0x28, 0x28)
GRAY        = RGBColor(0x90, 0x9A, 0xA0)
GREEN_CARD  = RGBColor(0xE8, 0xF5, 0xE9)
AMBER_CARD  = RGBColor(0xFF, 0xF3, 0xE0)
RED_CARD    = RGBColor(0xFF, 0xEB, 0xEE)
ACCENT_CARD = RGBColor(0xE3, 0xF2, 0xFD)
GRAY_CARD   = RGBColor(0xF5, 0xF5, 0xF5)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return s


def txb(sl, text, x, y, w, h,
        size=20, bold=False, color=TEXT,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = sl.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def bullet_box(sl, items, x, y, w, h, size=17, color=TEXT):
    tb = sl.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for text, level in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6 if level == 0 else 2)
        bullet = "•" if level == 0 else "◦"
        indent = "    " * level
        run = p.add_run()
        run.text = f"{indent}{bullet}  {text}"
        run.font.size = Pt(size if level == 0 else size - 2)
        run.font.color.rgb = color if level == 0 else GRAY
    return tb


def divider(sl, y=0.95, color=ACCENT):
    line = sl.shapes.add_connector(
        1, Inches(0.5), Inches(y), Inches(12.83), Inches(y))
    line.line.color.rgb = color
    line.line.width = Pt(1.2)


def accent_bar(sl, color=ACCENT, w=0.07):
    box = sl.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(w), H)
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()


def section_tag(sl, text, x=0.5, y=0.13, color=ACCENT):
    txb(sl, text, x, y, 6, 0.42, size=11, bold=True, color=color)


def tint(c, factor=0.88):
    h = str(c)
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return RGBColor(
        int(r + (255 - r) * factor),
        int(g + (255 - g) * factor),
        int(b + (255 - b) * factor))


# ══════════════════════════════════════════════════════════════
#  SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════
s = slide()
accent_bar(s, ACCENT)

txb(s, "Edit Entailment Learning",
    0.6, 1.8, 11.5, 1.2, size=44, bold=True, color=DARK)
txb(s, "Detecting Scope Creep in LLM-Generated Code Patches",
    0.6, 3.0, 11.5, 0.8, size=24, color=ACCENT)
txb(s, "When a requirement changes — which edits does the codebase truly need?",
    0.6, 3.85, 11.0, 0.6, size=17, color=SUB)
txb(s, "ICLR 2026 Submission",
    0.6, 6.6, 5, 0.5, size=13, color=GRAY)


# ══════════════════════════════════════════════════════════════
#  SLIDE 2 — Motivation
# ══════════════════════════════════════════════════════════════
s = slide()
accent_bar(s, RED)
section_tag(s, "MOTIVATION", color=RED)
divider(s, 0.75, RED)

txb(s, "Do you trust every hunk in an AI-generated PR?",
    0.6, 0.85, 12, 0.65, size=26, bold=True, color=DARK)
txb(s, "Scenario: an AI coding agent submits a PR with 5 hunks",
    0.6, 1.65, 12, 0.45, size=17, color=SUB)

colors_h   = [GREEN, GREEN, GREEN, AMBER, RED]
card_colors = [GREEN_CARD, GREEN_CARD, GREEN_CARD, AMBER_CARD, RED_CARD]
labels     = ["hunk 1\n✓ needed", "hunk 2\n✓ needed",
              "hunk 3\n✓ needed", "hunk 4\n? maybe", "hunk 5\n✗ extra"]
for i, (c, cc, lbl) in enumerate(zip(colors_h, card_colors, labels)):
    bx = s.shapes.add_shape(1,
        Inches(0.6 + i * 2.4), Inches(2.25), Inches(2.1), Inches(1.4))
    bx.fill.solid(); bx.fill.fore_color.rgb = cc
    bx.line.color.rgb = c; bx.line.width = Pt(2)
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = lbl
    run.font.size = Pt(14); run.font.bold = True; run.font.color.rgb = c

txb(s, "Scope Creep: AI generates edits beyond the scope of the requirement",
    0.6, 3.9, 12, 0.45, size=16, color=AMBER, bold=True)
bullet_box(s, [
    ("Existing tools check req. similarity OR test coverage — cannot identify hunks 4 and 5", 0),
    ("Cost: silent codebase pollution and heavier code review burden", 0),
], 0.6, 4.45, 12.2, 1.5, size=15, color=SUB)


# ══════════════════════════════════════════════════════════════
#  SLIDE 3 — The Hard Problem: T2
# ══════════════════════════════════════════════════════════════
s = slide()
accent_bar(s, AMBER)
section_tag(s, "THE HARD PROBLEM", color=AMBER)
divider(s, 0.75, AMBER)

txb(s, "Invisible Necessity: T2 is the blind spot of all existing methods",
    0.6, 0.85, 12.5, 0.6, size=24, bold=True, color=DARK)

tiers = [
    ("T1", "Gold + test-covered",
     "Req-driven, behavior visible",
     "All methods find it", GREEN, GREEN_CARD),
    ("T2", "Gold + no test coverage",
     "Structurally necessary, but tests miss it",
     "Systematically missed by all tools  ⚠️", AMBER, AMBER_CARD),
    ("T3", "LLM-generated (non-gold)",
     "Beyond requirement scope",
     "Scope creep — must be filtered", RED, RED_CARD),
]
for i, (tag, sub, desc, note, c, cc) in enumerate(tiers):
    x = 0.5 + i * 4.2
    bx = s.shapes.add_shape(1, Inches(x), Inches(1.65), Inches(3.9), Inches(4.6))
    bx.fill.solid(); bx.fill.fore_color.rgb = cc
    bx.line.color.rgb = c; bx.line.width = Pt(2)
    txb(s, tag,  x+0.15, 1.8,  3.6, 0.65, size=32, bold=True, color=c)
    txb(s, sub,  x+0.15, 2.55, 3.6, 0.5,  size=15, bold=True, color=TEXT)
    txb(s, desc, x+0.15, 3.1,  3.6, 0.85, size=13, color=SUB)
    txb(s, note, x+0.15, 4.65, 3.6, 0.85, size=13, color=c, bold=True)


# ══════════════════════════════════════════════════════════════
#  SLIDE 4 — Key Insight
# ══════════════════════════════════════════════════════════════
s = slide()
accent_bar(s, ACCENT)
section_tag(s, "KEY INSIGHT", color=ACCENT)
divider(s, 0.75, ACCENT)

txb(s, "Edit necessity exists on a continuum",
    0.6, 0.9, 12, 0.65, size=30, bold=True, color=DARK)

for i in range(100):
    t = i / 99
    rv = int(0xC6 * t + 0x2E * (1 - t))
    gv = int(0x28 * t + 0x7D * (1 - t))
    bv = int(0x28 * t + 0x32 * (1 - t))
    bar = s.shapes.add_shape(1,
        Inches(0.6 + i * 0.122), Inches(1.85), Inches(0.13), Inches(0.45))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(rv, gv, bv)
    bar.line.fill.background()

txb(s, "← Necessary (gold patch)", 0.6, 2.38, 5, 0.38, size=12, color=GREEN)
txb(s, "Scope Creep →", 8.4, 2.38, 4, 0.38, size=12, color=RED,
    align=PP_ALIGN.RIGHT)

questions = [
    (ACCENT, "①  Is it described by the requirement text?",
     "req similarity"),
    (GREEN,  "②  Is it constrained by fail-to-pass tests?",
     "test coverage"),
    (AMBER,  "③  Is it induced by the original code structure?",
     "structural context  ← NEW signal"),
]
for i, (c, q, label) in enumerate(questions):
    txb(s, q,     0.6,  2.95 + i * 0.9, 9.8, 0.5, size=17, color=c, bold=True)
    txb(s, label, 10.5, 2.95 + i * 0.9, 2.7, 0.5, size=12, color=c)

txb(s, "Our approach: jointly model all three signals — learn a geometric necessity gradient",
    0.6, 5.85, 12.2, 0.6, size=17, bold=True, color=ACCENT,
    align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  SLIDE 5 — Method
# ══════════════════════════════════════════════════════════════
s = slide()
accent_bar(s, ACCENT)
section_tag(s, "METHOD", color=ACCENT)
divider(s, 0.75, ACCENT)

txb(s, "Four-Entity Joint Embedding Space",
    0.6, 0.85, 11, 0.6, size=28, bold=True, color=DARK)

entities = [
    ("<REQ>",  "requirement text",  RGBColor(103,  58, 183)),
    ("<TEST>", "test fn. body",     RGBColor( 21, 101, 192)),
    ("<ORIG>", "source unit",       RGBColor(230,  81,   0)),
    ("<HUNK>", "diff hunk",         RGBColor(  0, 121, 107)),
]
ent_cards = [
    RGBColor(0xED, 0xE7, 0xF6),
    RGBColor(0xE3, 0xF2, 0xFD),
    RGBColor(0xFF, 0xF3, 0xE0),
    RGBColor(0xE0, 0xF2, 0xF1),
]
for i, ((tag, desc, c), cc) in enumerate(zip(entities, ent_cards)):
    x = 0.5 + i * 3.1
    bx = s.shapes.add_shape(1, Inches(x), Inches(1.65), Inches(2.85), Inches(1.05))
    bx.fill.solid(); bx.fill.fore_color.rgb = cc
    bx.line.color.rgb = c; bx.line.width = Pt(1.5)
    txb(s, tag,  x+0.1, 1.72, 2.6, 0.42, size=14, bold=True, color=c)
    txb(s, desc, x+0.1, 2.15, 2.6, 0.38, size=12, color=SUB)

txb(s, "↓   Shared UniXCoder Encoder   ↓",
    3.5, 2.85, 6.5, 0.44, size=15, color=GRAY, align=PP_ALIGN.CENTER)

txb(s, "Training: MultiPair InfoNCE — 4 positive pair types pulled together; in-batch negatives pushed apart",
    0.6, 3.42, 12.2, 0.44, size=14, color=SUB)

pairs = [
    ("(REQ, HUNK_gold)", RGBColor(103, 58, 183)),
    ("(REQ, TEST)",      RGBColor( 21,101, 192)),
    ("(ORIG, HUNK)",     RGBColor(230,  81,   0)),
    ("(REQ, ORIG)",      RGBColor(  0,121, 107)),
]
pair_cards = [
    RGBColor(0xED, 0xE7, 0xF6),
    RGBColor(0xE3, 0xF2, 0xFD),
    RGBColor(0xFF, 0xF3, 0xE0),
    RGBColor(0xE0, 0xF2, 0xF1),
]
for i, ((pt, c), cc) in enumerate(zip(pairs, pair_cards)):
    bx = s.shapes.add_shape(1, Inches(0.5 + i*3.1), Inches(3.98),
                             Inches(2.85), Inches(0.58))
    bx.fill.solid(); bx.fill.fore_color.rgb = cc
    bx.line.color.rgb = c; bx.line.width = Pt(1.2)
    txb(s, pt, 0.6+i*3.1, 4.06, 2.6, 0.46, size=13, color=c, bold=True)

txb(s, "Inference score:",
    0.6, 4.82, 3.5, 0.44, size=15, color=SUB)

tb = s.shapes.add_textbox(Inches(4.1), Inches(4.72), Inches(8.5), Inches(0.72))
tf = tb.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "s(h)  =  α·sim(h, r)   +   β·sim(h, t̅)   +   γ·sim(h, o̅)"
run.font.size = Pt(19); run.font.bold = True; run.font.color.rgb = DARK

txb(s, "α = 0.4    β = 0.3    γ = 0.3",
    4.1, 5.48, 8.5, 0.44, size=13, color=GRAY)


# ══════════════════════════════════════════════════════════════
#  SLIDE 6 — Learned Geometry
# ══════════════════════════════════════════════════════════════
s = slide()
accent_bar(s, ACCENT)
section_tag(s, "LEARNED GEOMETRY", color=ACCENT)
divider(s, 0.75, ACCENT)

txb(s, "Five Semantic Regions Emerge from Training",
    0.6, 0.85, 12, 0.6, size=26, bold=True, color=DARK)

ph = s.shapes.add_shape(1, Inches(0.5), Inches(1.5), Inches(6.8), Inches(4.8))
ph.fill.solid(); ph.fill.fore_color.rgb = GRAY_CARD
ph.line.color.rgb = GRAY; ph.line.width = Pt(1)
txb(s, "[ Insert: django_combined.png ]",
    0.5, 3.75, 6.8, 0.5, size=13, color=GRAY, align=PP_ALIGN.CENTER)

txb(s, "Region    Active Anchors          Interpretation",
    7.5, 1.5, 5.8, 0.4, size=11, bold=True, color=SUB)

regions = [
    ("R1", "REQ + TEST + ORIG", "Direct fix — test-visible",           GREEN),
    ("R2", "REQ + ORIG",        "Necessary but test-invisible",         ACCENT),
    ("R3", "REQ only",          "New code, no prior structure",         SUB),
    ("R4", "ORIG only",         "Structurally induced  ⚡",         AMBER),
    ("R5", "None",              "Scope Creep",                          RED),
]
for i, (r, anchors, meaning, c) in enumerate(regions):
    y = 2.05 + i * 0.88
    card_c = tint(c, 0.9)
    bx = s.shapes.add_shape(1, Inches(7.4), Inches(y - 0.05),
                             Inches(5.8), Inches(0.78))
    bx.fill.solid(); bx.fill.fore_color.rgb = card_c
    bx.line.color.rgb = c; bx.line.width = Pt(1)
    txb(s, r,       7.5,  y, 0.7, 0.65, size=14, bold=True, color=c)
    txb(s, anchors, 8.2,  y, 2.3, 0.65, size=12, color=TEXT)
    txb(s, meaning, 10.6, y, 2.5, 0.65, size=12, color=c)

txb(s, "R4 = the core value of ORIG — structurally necessary edits the requirement never mentions",
    0.5, 6.55, 12.5, 0.55, size=13, bold=True, color=AMBER,
    align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  SLIDE 7 — Results
# ══════════════════════════════════════════════════════════════
s = slide()
accent_bar(s, GREEN)
section_tag(s, "RESULTS", color=GREEN)
divider(s, 0.75, GREEN)

txb(s, "Cross-Repo Generalization: Learned Geometry, Not Memorization",
    0.6, 0.85, 12.5, 0.6, size=24, bold=True, color=DARK)

nums = [
    ("0.825", "nDCG@k",      "Overall ranking quality"),
    ("0.839", "T2-Recall",   "Invisible necessary edits"),
    ("5.8×",  "Silhouette↑", "Geometric separation"),
]
for i, (val, metric, desc) in enumerate(nums):
    x = 0.5 + i * 4.2
    bx = s.shapes.add_shape(1, Inches(x), Inches(1.65), Inches(3.8), Inches(1.85))
    bx.fill.solid(); bx.fill.fore_color.rgb = GREEN_CARD
    bx.line.color.rgb = GREEN; bx.line.width = Pt(2)
    txb(s, val,    x+0.15, 1.75, 3.5, 0.75, size=36, bold=True, color=GREEN)
    txb(s, metric, x+0.15, 2.48, 3.5, 0.38, size=15, bold=True, color=DARK)
    txb(s, desc,   x+0.15, 2.88, 3.5, 0.38, size=11, color=SUB)

txb(s, "Cross-repo validation across 12 Python repositories:",
    0.6, 3.75, 12, 0.44, size=17, bold=True, color=ACCENT)

bullet_box(s, [
    ("pallets (only 11 training instances): T2-Recall +0.07  — largest gain in the smallest repo", 0),
    ("mwaskom (22 instances) → +0.11   /   psf (44 instances) → +0.10", 0),
    ("Gains are NOT correlated with training size → transferable geometric structure, not memorization", 0),
], 0.6, 4.3, 12.2, 2.8, size=15, color=TEXT)


# ══════════════════════════════════════════════════════════════
#  SLIDE 8 — Ablation
# ══════════════════════════════════════════════════════════════
s = slide()
accent_bar(s, AMBER)
section_tag(s, "ABLATION", color=AMBER)
divider(s, 0.75, AMBER)

txb(s, "Every Entity Is Indispensable",
    0.6, 0.85, 12, 0.6, size=28, bold=True, color=DARK)

hbg = s.shapes.add_shape(1, Inches(0.45), Inches(1.65), Inches(12.45), Inches(0.48))
hbg.fill.solid(); hbg.fill.fore_color.rgb = ACCENT_CARD
hbg.line.fill.background()

col_x = [0.5, 5.1, 7.1, 9.1]
col_w = [4.4, 1.8, 1.8, 4.0]
headers = ["Model Variant", "nDCG@k", "T2-Recall", "Primary Impact"]
for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    txb(s, hdr, cx, 1.7, cw, 0.4, size=13, bold=True, color=ACCENT)

rows_data = [
    ("M2  (full model)",    "0.825 ✓", "0.839 ✓", "—"),
    ("w/o ORIG  (γ = 0)", "0.668 ↓", "0.721 ↓", "R4 region collapses"),
    ("w/o TEST  (β = 0)", "0.806",        "0.830",        "T1/T2 gradient blurs"),
    ("w/o REQ   (α = 0)", "0.789",        "0.755 ↓", "Semantic coverage drops"),
]
row_bgs = [GREEN_CARD, RED_CARD, AMBER_CARD, GRAY_CARD]

for i, (row, rbg) in enumerate(zip(rows_data, row_bgs)):
    y = 2.25 + i * 0.9
    rb = s.shapes.add_shape(1, Inches(0.45), Inches(y - 0.05),
                             Inches(12.45), Inches(0.84))
    rb.fill.solid(); rb.fill.fore_color.rgb = rbg
    rb.line.fill.background()
    for j, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
        if "✓" in cell:
            c = GREEN
        elif "↓" in cell:
            c = RED
        elif i == 0 and j == 0:
            c = GREEN
        elif j == 0:
            c = AMBER
        else:
            c = TEXT
        txb(s, cell, cx, y, cw, 0.78, size=14, color=c, bold=(j == 0))

txb(s, "ORIG is the critical signal: removing it drops T2-Recall by −0.118",
    0.6, 5.9, 12.2, 0.48, size=16, bold=True, color=AMBER,
    align=PP_ALIGN.CENTER)
txb(s, "ORIG is the only entity that captures structurally-induced necessity — without it, R4 collapses",
    0.6, 6.42, 12.2, 0.55, size=13, color=SUB,
    align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  SLIDE 9 — Conclusion
# ══════════════════════════════════════════════════════════════
s = slide()
accent_bar(s, ACCENT)
section_tag(s, "CONCLUSION", color=ACCENT)
divider(s, 0.75, ACCENT)

txb(s, "EEL Learns Geometric Structure, Not a Classifier",
    0.6, 0.85, 12, 0.6, size=28, bold=True, color=DARK)

bullet_box(s, [
    ("Edit necessity is a continuous gradient, expressible as learned geometry", 0),
    ("The ORIG entity unlocks T2 detection — invisible to all existing methods", 0),
    ("Cross-repo generalization confirms the geometry transfers across codebases", 0),
], 0.6, 1.6, 12.2, 1.9, size=17, color=TEXT)

txb(s, "Limitations",  0.6, 3.65, 5, 0.44, size=17, bold=True, color=AMBER)
bullet_box(s, [
    ("Automated labeling noise: some LLM-generated T3 hunks may be reasonable implementations", 0),
    ("Inference weights α / β / γ are fixed globally — no per-instance adaptation", 0),
], 0.6, 4.18, 12.2, 1.2, size=14, color=SUB)

txb(s, "Future Work", 0.6, 5.5, 5, 0.44, size=17, bold=True, color=ACCENT)
bullet_box(s, [
    ("Human-annotated soft labels → training on continuous necessity scores", 0),
    ("Region-aware curriculum learning: R1 vs R5 first, then R2 vs R4", 0),
    ("Transfer to PR review automation and real-time AI coding agent filtering", 0),
], 0.6, 6.02, 12.2, 1.3, size=14, color=TEXT)


# ══════════════════════════════════════════════════════════════
#  SLIDE 10 — Thank You
# ══════════════════════════════════════════════════════════════
s = slide()
accent_bar(s, ACCENT)

txb(s, "Thank You!",
    0.6, 2.0, 12.2, 1.2, size=52, bold=True, color=DARK,
    align=PP_ALIGN.CENTER)
txb(s, "Edit Entailment Learning",
    0.6, 3.4, 12.2, 0.6, size=22, color=ACCENT,
    align=PP_ALIGN.CENTER)
txb(s, "Q & A",
    0.6, 4.3, 12.2, 0.6, size=28, color=SUB,
    align=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────
out = "EEL_presentation.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
